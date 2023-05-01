from pptx import Presentation
from pptx.util import Inches


# Criando a apresentação
prs = Presentation()

# Criando o primeiro slide
slide = prs.slides.add_slide(prs.slide_layouts[0])

# Adicionando um título ao slide
title = slide.shapes.title
title.text = "Animais"

# Adicionando uma imagem de um cachorro ao slide
img_path = 'cachorro.jpg'
left = top = Inches(2)
pic = slide.shapes.add_picture(img_path, left, top)

# Criando o segundo slide
slide = prs.slides.add_slide(prs.slide_layouts[1])

# Adicionando um título ao slide
title = slide.shapes.title
title.text = "Gatos"

# Adicionando uma imagem de um gato ao slide
img_path = 'gato.jpg'
left = top = Inches(2)
pic = slide.shapes.add_picture(img_path, left, top)

# Salvando a apresentação
prs.save('animais.pptx')
